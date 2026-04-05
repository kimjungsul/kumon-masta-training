"""
make_pptx.py
실행하면 마스타_교사교육안.pptx 파일을 생성합니다.

사용법:
  python make_pptx.py

데이터 수정:
  아래 데이터 정의 영역만 수정하면 자동으로 슬라이드가 업데이트됩니다.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_FILE = "d:/Coding/claude/마스타_교사교육안.pptx"

# ──────────────────────────────────────────
# 데이터 정의 (이 부분만 수정하면 됩니다)
# ──────────────────────────────────────────

META = {
    "title": "마.스.타 활성화 교사 교육안",
    "period": "2026년 4월 ~ 6월 (3개월)",
    "goal": "신규 47명 확보 / 이탈률 50% 감소",
    "target": "전체 교사 11명",
}

CURRENT_STATUS = [
    ("전체 관리회원",     "483명"),
    ("마.스.타 회원",     "64명 (13.3%)"),
    ("스마트 구몬 회원",  "231명 (47.8%)"),
    ("미가입 대상",       "약 188명"),
    ("월 신규",           "약 10명"),
    ("월 이탈",           "약 13명"),
]

TEACHERS = [
    # (이름, 회원수, 마스타수, 점유율, 잔여대상, 2분기목표)
    ("김은숙", 103, 14, "13.6%", 89, "6~7명"),
    ("최승희",  68,  7, "10.3%", 61, "4~5명"),
    ("원희선",  50,  2,  "4.0%", 48, "3~4명"),
    ("박양희",  50,  2,  "4.0%", 48, "3~4명"),
    ("이경옥",  47,  3,  "6.4%", 44, "3명"),
    ("김주연",  51, 14, "27.5%", 37, "4명"),
    ("이경이",  35,  1,  "2.9%", 34, "4명"),
    ("이현주",  29,  2,  "6.9%", 27, "3명"),
    ("이주아",  19,  1,  "5.3%", 18, "2명"),
    ("신인숙",  16,  3, "18.8%", 13, "1명"),
    ("김태연",  18,  1,  "5.6%", 17, "1명"),
]

EDUCATION_SESSIONS = [
    ("1회차", "가치 전달 화법",       "4월 첫째 주"),
    ("2회차", "이탈 방어",            "4월 셋째 주"),
    ("3회차", "목표 관리 & 실전 점검","5월 첫째 주"),
]

MONTHLY_TARGETS = [
    # (월, 마스타목표, 퇴회율목표, 퇴회수목표)
    ("4월", "+11명", "5.2% 이하", "81.5명↓"),
    ("5월", "+12명", "5.7% 이하", "89명↓"),
    ("6월", "+24명", "3.8% 이하", "60명↓"),
]

SUCCESS_CRITERIA = [
    ("마.스.타 회원 수", "64명",  "93명 이상"),
    ("마.스.타 점유율", "13.3%", "약 19%"),
    ("월 이탈",         "13명",  "6명 이하"),
    ("월 신규",         "10명",  "16명 이상"),
]

# 전략 2 — 퇴회율 안정화
CHURN_MONTHLY = [
    # (월, 퇴회율목표, 퇴회수목표, 핵심액션)
    ("4월", "5.2%", "81.5명↓", "이탈 조기 경보 체계 구축"),
    ("5월", "5.7%", "89명↓",   "재약정 선제 상담 병행"),
    ("6월", "3.8%", "60명↓",   "5월 말 전 회원 선제 상담 완료"),
]

CHURN_SIGNALS = [
    ("숙제 미완료 2회 연속",  "분량 30% 축소 → 성공 경험 후 복원"),
    ('"어려워요" 반복',       "1단계 하향, 기초 보강 후 복귀"),
    ('"그만둘까 해요"',       "48시간 내 전화 상담, 조정 옵션 제시"),
    ("시간 조율 어려움",      "2개 타임슬롯 제공 또는 교사 간 교환 조율"),
]

# 전략 3 — SKN 재약정률
RECONTRACT_MONTHLY = [
    # (월, 대상, 목표계약, 목표율, 비고)
    ("4월", "7명",  "6명",  "85.7%", "만료 1개월 전 선 상담"),
    ("5월", "8명",  "7명",  "87.5%", "고맙구몬 성장노트 활용"),
    ("6월", "6명",  "6명",  "100%",  "관리자 동행 상담 필수"),
    ("합계","21명", "19명", "90.5%", ""),
]

# SKN 50% 달성 분석
SKN_ANALYSIS = {
    "현재_총회원":    483,
    "현재_SKN":       231,
    "현재_비율":      "47.8%",
    "2분기_순증":     38.5,
    "2분기_SKN계약":  40,
    "2분기_SKN해지":  22,
}
SKN_ANALYSIS["말_총회원"]  = SKN_ANALYSIS["현재_총회원"] + SKN_ANALYSIS["2분기_순증"]
SKN_ANALYSIS["말_SKN"]     = SKN_ANALYSIS["현재_SKN"] + SKN_ANALYSIS["2분기_SKN계약"] - SKN_ANALYSIS["2분기_SKN해지"]
SKN_ANALYSIS["말_비율"]    = round(SKN_ANALYSIS["말_SKN"] / SKN_ANALYSIS["말_총회원"] * 100, 1)
SKN_ANALYSIS["50%_필요"]   = round(SKN_ANALYSIS["말_총회원"] * 0.5)
SKN_ANALYSIS["추가_필요"]  = SKN_ANALYSIS["50%_필요"] - SKN_ANALYSIS["말_SKN"]

# 전략 4 — 교사 순증
TEACHER_RECRUIT = [
    # (월, 목표순증, 방법)
    ("4월", "+1명", "청년교사 공고 등록"),
    ("5월", "+1명", "기존 교사 추천 캠페인"),
    ("6월", "+3명", "하반기 대비 집중 채용"),
    ("합계","+4명 (31→35명)", ""),
]

# ──────────────────────────────────────────
# 스타일 설정
# ──────────────────────────────────────────

TITLE_COLOR  = RGBColor(0x1A, 0x56, 0xDB)
BG_COLOR     = RGBColor(0xF8, 0xF9, 0xFF)
TEXT_COLOR   = RGBColor(0x1F, 0x29, 0x37)
SUB_COLOR    = RGBColor(0x6B, 0x72, 0x80)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ODD      = RGBColor(0xE8, 0xF0, 0xFF)
HIGHLIGHT    = RGBColor(0xFF, 0xED, 0xCC)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ──────────────────────────────────────────
# 헬퍼 함수
# ──────────────────────────────────────────

def bg(slide, color=BG_COLOR):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def bar(slide):
    s = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    s.fill.solid(); s.fill.fore_color.rgb = TITLE_COLOR
    s.line.fill.background()

def textbox(slide, text, l, t, w, h, size=18, bold=False,
            color=TEXT_COLOR, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def title_bar(slide, text):
    bar(slide)
    W = prs.slide_width
    textbox(slide, text, Inches(0.5), Inches(0.15), W - Inches(1), Inches(1),
            size=26, bold=True, color=TITLE_COLOR)

def make_table(slide, headers, rows, top=Inches(1.3), highlight_last=False):
    W = prs.slide_width
    cols = len(headers)
    row_h = 0.45
    tbl = slide.shapes.add_table(
        len(rows) + 1, cols,
        Inches(0.5), top,
        W - Inches(1), Inches(row_h * (len(rows) + 1))
    ).table

    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = TITLE_COLOR
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold  = True
        p.runs[0].font.size  = Pt(13)
        p.runs[0].font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(rows):
        is_last = (r == len(rows) - 1)
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            if highlight_last and is_last:
                cell.fill.solid(); cell.fill.fore_color.rgb = HIGHLIGHT
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ROW_ODD
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                p.runs[0].font.size  = Pt(12)
                p.runs[0].font.color.rgb = TEXT_COLOR

def bullets(slide, items, top=Inches(1.3)):
    W = prs.slide_width
    y = top
    for item in items:
        is_sub  = item.startswith("  ")
        is_note = item.startswith("※")
        is_head = item.startswith("▶")
        txt = item.strip()

        if txt == "":
            y += Inches(0.2)
            continue

        indent = Inches(1.0) if is_sub else Inches(0.5)
        fs     = 15 if is_sub else (14 if is_note else 17)
        col    = SUB_COLOR if is_note else (TITLE_COLOR if is_head else TEXT_COLOR)
        bold   = is_head
        prefix = "   · " if is_sub else ("" if (is_note or is_head) else "• ")

        textbox(slide, prefix + txt, indent, y,
                W - indent - Inches(0.4), Inches(0.55),
                size=fs, bold=bold, color=col)
        y += Inches(0.48) if is_sub else Inches(0.52)

# ──────────────────────────────────────────
# 슬라이드 생성
# ──────────────────────────────────────────

# 1. 표지
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0x1A, 0x56, 0xDB))
W, H = prs.slide_width, prs.slide_height
textbox(s, META["title"], Inches(1), Inches(1.8), W - Inches(2), Inches(1.6),
        size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
sub = f"기간: {META['period']}\n목표: {META['goal']}\n대상: {META['target']}"
textbox(s, sub, Inches(1), Inches(3.7), W - Inches(2), Inches(2),
        size=20, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)

# 2. 현황 수치
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "현황 분석 — 핵심 수치")
items = [f"{k}  {v}" for k, v in CURRENT_STATUS]
items += ["", "※ 월 신규(10명) < 월 이탈(13명) → 순감 진행 중"]
bullets(s, items)

# 3. 교사 분류
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "교사 분류")
bullets(s, [
    "▶ 우수 교사 (벤치마크 대상)",
    "  김주연 — 점유율 27.5% (전체 1위), 권유 노하우 공유 대상",
    "  김은숙 — 마스타 14명 최다 보유, 대규모 관리 역량 우수",
    "",
    "▶ 집중 육성 대상 (대상 풀 多, 점유율 低)",
    "  이경옥 47명 중 3명(6.4%) / 원희선 50명 중 2명(4.0%)",
    "  박양희 50명 중 2명(4.0%) / 최승희 68명 중 7명(10.3%)",
    "  이경이 35명 중 1명(2.9%) ← 전환 가능성 최고",
    "",
    "▶ 소규모 교사",
    "  신인숙(16명) / 이주아(19명) / 김태연(18명) → 현실적 목표 유지",
])

# 4. 이탈 3대 사유
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "이탈 3대 사유")
bullets(s, [
    "1.  월 4만원 추가 비용 대비 가치를 못 느낌",
    "2.  추가 학습량 부담으로 포기",
    "3.  교사와 시간 조율 어려움",
])

# 5. 교육 구성
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "교육 구성 (총 3회차)")
make_table(s, ["회차", "주제", "일정"],
           [[a, b, c] for a, b, c in EDUCATION_SESSIONS])

# 6. 1회차 화법
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "1회차 — 비용 저항 극복 화법")
make_table(s,
    ["학부모 반응", "전환 화법"],
    [
        ['"4만원이 부담돼요"',       '주 1회 기준 회당 1만원, 학원 보충수업보다 낮으면서 1:1 맞춤 지도'],
        ['"교재만 더 푸는 거 아닌가요?"', 'OO의 약한 부분만 집중 보강하는 구조, 진단 결과 기반 설명'],
        ['"효과를 모르겠어요"',      '3개월 전 시작한 OO 회원 진단평가 점수 데이터로 제시'],
    ]
)

# 7. 권유 타이밍
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "1회차 — 권유 황금 타이밍")
bullets(s, [
    "진단평가 후 결과 공유 시  →  마.스.타 효과 연결",
    "학부모 상담 중 학습 고민 언급 시  →  해결책으로 자연스럽게 제안",
    "스마트 구몬 데이터 리뷰 시  →  약점 영역 기반 제안",
    "",
    "※ 실습: 교사 2인 1조 롤플레이 (3가지 거절 상황 대응 연습)",
])

# 8. 이탈 방어
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "2회차 — 이탈 신호별 교사 조치")
make_table(s,
    ["이탈 신호", "교사 조치"],
    [
        ["숙제 미완료 2회 연속",    "분량 30% 축소 → 성공 경험 후 복원"],
        ['"어려워요" 반복',         "1단계 하향, 기초 보강 후 복귀"],
        ['"그만둘까 해요"',         "48시간 내 전화 상담, 조정 옵션 제시"],
    ]
)
textbox(s, "※ 목표: 월 이탈 13명 → 6명 이하", Inches(0.5), Inches(5.8),
        prs.slide_width - Inches(1), Inches(0.5), size=14, color=SUB_COLOR)

# 9. 시간 조율
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "2회차 — 시간 조율 개선안")
bullets(s, [
    "개선안 1  주 중 2개 타임슬롯 제공 (회원 선택)",
    "개선안 2  비대면 보충 지도 병행 (스마트 구몬 연계)",
    "개선안 3  교사 간 회원 시간대 교환 조율",
])

# 10. 교사별 목표 표
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "3회차 — 교사별 목표")
rows = [(n, f"{rem}명", g) for n, _, _, _, rem, g in TEACHERS]
rows.append(["합계", "436명", "47명"])
make_table(s, ["교사명", "잔여 대상", "2분기 목표"], rows, highlight_last=True)

# 11. KPI
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "3회차 — KPI & 보고 체계")
bullets(s, [
    "▶ 월간 KPI",
    "  마.스.타 권유 상담  →  월 8회 이상 (소규모 교사 4회)",
    "  이탈 신호 조기 상담  →  발생 48시간 내 100%",
    "  기존 회원 이탈  →  0명 목표",
    "",
    "▶ 보고 체계",
    "  매주 금요일  교사 → 관리자 (신규 상담·가입·이탈 위기)",
    "  매월 첫째 주  전체 교사 월간 실적 리뷰 미팅",
])

# 12. 2분기 전략 1 개요
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "2분기 전략 1 — 마.스타 순증 +47 달성")
bullets(s, ["※ 현황: 1분기 -4.5 → 2분기 목표 +47  (최우선 과제)", ""])
make_table(s, ["월", "마.스타 목표", "퇴회율 목표", "퇴회 수 목표"],
           [[m, t, r, n] for m, t, r, n in MONTHLY_TARGETS],
           top=Inches(2.3))
textbox(s, "※ 핵심 전제: 마.스타 퇴회 억제 없이는 순증 달성 불가",
        Inches(0.5), Inches(5.8), prs.slide_width - Inches(1), Inches(0.5),
        size=14, color=SUB_COLOR)

# 13. 대형 교사 집중 공략
priority = [t for t in TEACHERS if t[4] >= 40]
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "전략 1-① 대형 교사 집중 공략")
rows = [(n, f"{m}명", p, f"{rem}명", g) for n, _, m, p, rem, g in priority]
make_table(s, ["교사명", "현재 마.스타", "점유율", "잔여 대상", "2분기 목표"], rows)
bullets(s, [
    "",
    "※ 김은숙·최승희: 잔여 대상 최다 → 2분기 순증 핵심 동력",
    "※ 원희선·박양희: 점유율 4.0% 최하위 → 진단평가 기반 전환 상담 즉시 투입",
], top=Inches(4.8))

# 14. 우수 교사 노하우 전파
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "전략 1-② 우수 교사 노하우 전파")
bullets(s, [
    "▶ 김주연 교사  (점유율 27.5% — 전체 1위)",
    "  4월 교사 미팅에서 권유 화법·성공 사례 공유",
    "  저점유율 교사에게 스크립트 직접 적용",
    "",
    "▶ 김은숙 교사  (마.스타 14명 — 최다 보유)",
    "  대규모 회원 관리 시 권유 타이밍·노하우 공유",
    "",
    "※ 적용 대상: 원희선·박양희·이경이 (점유율 2.9~4.0% 최하위 3인)",
])

# 15. 스마트구몬 전환 + 소규모 교사
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "전략 1-③④ 전환 집중 & 소규모 교사 관리")
bullets(s, [
    "▶ ③ 스마트구몬 → 마.스타 전환 집중",
    "  SKN 계약 회원 중 마.스타 미가입자 교사별 목록화",
    "  이경이 교사 — 점유율 2.9%, 잔여 34명 → 관리자 동행 상담 우선 투입",
    "",
    "▶ ④ 소규모 교사 — 현실적 목표 유지",
    "  신인숙(잔여 13명) / 이주아(18명) / 김태연(17명)",
    "  2분기 각 1명 목표 / 권유 상담 월 4회 이상 이행 점검",
    "",
    "※ 전체 합계: 47명 = 대형 교사 20명 + 중형 10명 + 저점유율 전환 17명",
])

# ── 전략 2 ──────────────────────────────────────────

# 16. 전략 2 개요
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "2분기 전략 2 — 퇴회율 안정화")
bullets(s, ["※ 현황: 1분기 6.0% (3월 8.8% 최고점) → 2분기 목표 4.9%", ""])
make_table(s, ["월", "퇴회율 목표", "퇴회 수 목표", "핵심 액션"],
           [[m, r, n, a] for m, r, n, a in CHURN_MONTHLY],
           top=Inches(2.3))
textbox(s, "※ 3월 퇴회율 8.8% 재발 시 전체 목표 붕괴 — 조기 경보 체계 필수",
        Inches(0.5), Inches(5.8), prs.slide_width - Inches(1), Inches(0.5),
        size=14, color=SUB_COLOR)

# 17. 이탈 신호 대응
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "전략 2 — 이탈 신호별 즉시 대응")
make_table(s, ["이탈 신호", "교사 조치"],
           [[sig, act] for sig, act in CHURN_SIGNALS])
bullets(s, [
    "",
    "▶ 6월 집중 방어 (목표 퇴회율 3.8% — 분기 최저)",
    "  5월 말 전 회원 선제 상담 완료",
    "  여름방학 연계 동기부여 (학습 목표 재설정)",
    "  미달 위험 회원 관리자 직접 동행",
], top=Inches(4.0))

# ── 전략 3 ──────────────────────────────────────────

# 18. 전략 3 개요
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "2분기 전략 3 — SKN 재약정률 90.5% 달성")
bullets(s, ["※ 현황: 1분기 73.7% (1월 50% 저점) → 2분기 목표 90.5%", ""])
make_table(s, ["월", "재약정 대상", "목표 계약", "목표율", "비고"],
           [[m, d, c, r, n] for m, d, c, r, n in RECONTRACT_MONTHLY],
           top=Inches(2.3), highlight_last=True)

# 19. 재약정 실행 액션
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "전략 3 — 재약정 실행 액션")
bullets(s, [
    "▶ 재약정 대상 사전 관리",
    "  만료 1개월 전 교사 선 상담 완료 → 관리자 확인",
    "  고맙구몬 스티커판·성장노트로 학습 성과 가시화",
    "  진단평가 결과 공유로 마.스.타 연계 효과 입증",
    "",
    "▶ 6월 100% 달성 (대상 6명 전원 계약 필수)",
    "  관리자 직접 동행 상담",
    "  체험 후 관리, 외국어·나무시리즈 성과 공유",
    "",
    "※ 재약정 미달 교사 — 5월 중간 점검 시 조기 식별 및 코칭",
])

# ── 전략 4 ──────────────────────────────────────────

# 20. 전략 4 개요
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "2분기 전략 4 — 교사 순증 +4명 (31→35명)")
make_table(s, ["월", "목표 순증", "주요 방법"],
           [[m, t, a] for m, t, a in TEACHER_RECRUIT],
           highlight_last=True)
bullets(s, [
    "",
    "▶ 청년교사제도 즉시 공고 등록 (4월)",
    "▶ 기존 교사 추천 시 인센티브 제공 (5월)",
    "▶ 6월 집중 채용 — 하반기 회원 증가 대비 선제 조직 확대",
    "※ 20과목 이하 교사 3명 → 지역 정리 및 분리 검토 병행",
], top=Inches(4.3))

# ── SKN 50% 분석 ──────────────────────────────────────────

# 21. SKN 50% 달성 분석
a = SKN_ANALYSIS
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "SKN 스마트구몬 50% 달성 분석")
make_table(s,
    ["항목", "현재", "2분기 말 예상"],
    [
        ["총 회원",     f"{a['현재_총회원']}명",  f"{int(a['말_총회원'])}명"],
        ["SKN 회원",    f"{a['현재_SKN']}명",     f"{int(a['말_SKN'])}명"],
        ["SKN 비율",    a["현재_비율"],           f"{a['말_비율']}%"],
        ["50% 달성 필요", "—",                   f"{int(a['50%_필요'])}명"],
    ]
)
bullets(s, [
    "",
    f"※ 현재 계획 달성 시 SKN {a['말_비율']}% → 50% 달성까지 {int(a['추가_필요'])}명 부족",
    f"※ 50% 달성을 위한 SKN 순증 목표: +18명 → +{18 + int(a['추가_필요'])}명으로 상향 필요",
    "※ 재약정률 90.5% 미달 시 SKN 해지 추가 발생 → 비율 더 하락",
], top=Inches(4.6))

# 22. SKN 50% 달성 전략
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "SKN 50% 달성 추가 전략 (+12명)")
bullets(s, [
    "▶ 신규 입회 시 SKN 동시 가입 유도 (월 입회 목표의 30% 이상)",
    "  4월 83명 입회 목표 중 SKN 동시 가입 25명 이상",
    "  입회 상담 시 SKN 체험 패키지 기본 제안",
    "",
    "▶ 기존 스마트구몬 미가입 회원 전환",
    "  전체 미가입 대상 약 188명 중 교사별 5명 선정 → 상담",
    "  진단평가 결과 기반 SKN 필요성 설명",
    "",
    "▶ 재약정률 90.5% 반드시 달성 (해지 방어)",
    "  재약정 실패 1명 = SKN 순증 목표 1명 손실과 동일",
    "  만료 전 관리자 동행 상담으로 해지 최소화",
])

# ── 종합 로드맵 ──────────────────────────────────────────

# 23. 월별 종합 로드맵
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "2분기 월별 종합 실행 로드맵")
make_table(s,
    ["월", "전략 1 (마.스타)", "전략 2 (퇴회)", "전략 3 (재약정)", "전략 4 (교사)"],
    [
        ["4월", "전환 상담 집중\n+11명 목표", "조기 경보 체계\n5.2% 이하", "선 상담 완료\n85.7% 목표", "청년교사 공고\n+1명"],
        ["5월", "권유 본격화\n+12명 목표",  "재약정 병행\n5.7% 이하",  "성장노트 활용\n87.5% 목표", "추천 캠페인\n+1명"],
        ["6월", "상반기 스퍼트\n+24명 목표", "전원 선제 상담\n3.8% 이하", "동행 상담\n100% 목표",   "집중 채용\n+3명"],
    ]
)

# 24. 성공 기준
s = prs.slides.add_slide(blank); bg(s); title_bar(s, "성공 기준")
rows = [[a, b, c] for a, b, c in SUCCESS_CRITERIA]
make_table(s, ["지표", "현재", "3개월 목표"], rows)

# ──────────────────────────────────────────
# 저장
# ──────────────────────────────────────────

prs.save(OUTPUT_FILE)
print(f"완료: {OUTPUT_FILE}")
print(f"총 슬라이드 수: {len(prs.slides)}장")
